"""
Document Processor
==================
Extracts text from various document formats (PDF, DOCX, PPTX, HTML, TXT)
and splits them into semantic chunks for embedding.

Math-aware pipeline:
  1. pymupdf4llm for structured markdown extraction (tables, headings preserved)
  2. Unicode math symbol normalization for better embedding retrieval
  3. Formula-aware chunking (equation blocks protected from mid-split)
  4. Dual-text storage: original text for LLM, normalized text for embedding
"""

import io
import logging
import re
import time
from dataclasses import dataclass, field
from pathlib import Path

from core import config

logger = logging.getLogger(__name__)

# OCR support (optional — graceful degrade if not installed)
try:
    import pytesseract
    from PIL import Image

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


# ─── Math Symbol Normalization Map ────────────────────────────────────────────

MATH_SYMBOL_MAP = {
    # Greek lowercase
    "\u03b1": "alpha",
    "\u03b2": "beta",
    "\u03b3": "gamma",
    "\u03b4": "delta",
    "\u03b5": "epsilon",
    "\u03b6": "zeta",
    "\u03b7": "eta",
    "\u03b8": "theta",
    "\u03b9": "iota",
    "\u03ba": "kappa",
    "\u03bb": "lambda",
    "\u03bc": "mu",
    "\u03bd": "nu",
    "\u03be": "xi",
    "\u03c0": "pi",
    "\u03c1": "rho",
    "\u03c3": "sigma",
    "\u03c4": "tau",
    "\u03c5": "upsilon",
    "\u03c6": "phi",
    "\u03c7": "chi",
    "\u03c8": "psi",
    "\u03c9": "omega",
    # Greek uppercase
    "\u0393": "Gamma",
    "\u0394": "Delta",
    "\u0398": "Theta",
    "\u039b": "Lambda",
    "\u039e": "Xi",
    "\u03a0": "Pi",
    "\u03a3": "Sigma",
    "\u03a6": "Phi",
    "\u03a8": "Psi",
    "\u03a9": "Omega",
    # Operators
    "\u222b": " integral ",
    "\u222c": " double integral ",
    "\u222d": " triple integral ",
    "\u2211": " sum ",
    "\u220f": " product ",
    "\u221a": " sqrt ",
    "\u221b": " cube root ",
    "\u2202": " partial ",
    "\u2207": " nabla ",
    "\u221e": " infinity ",
    # Relations
    "\u2260": " != ",
    "\u2264": " <= ",
    "\u2265": " >= ",
    "\u2248": " approximately ",
    "\u2261": " equivalent ",
    "\u221d": " proportional ",
    # Set theory
    "\u2208": " in ",
    "\u2209": " not in ",
    "\u2282": " subset ",
    "\u2283": " superset ",
    "\u2286": " subset_eq ",
    "\u2287": " superset_eq ",
    "\u222a": " union ",
    "\u2229": " intersection ",
    "\u2205": " empty set ",
    # Logic
    "\u2227": " and ",
    "\u2228": " or ",
    "\u00ac": " not ",
    "\u2200": " for all ",
    "\u2203": " there exists ",
    "\u21d2": " implies ",
    "\u21d4": " iff ",
    # Arrows
    "\u2192": " -> ",
    "\u2190": " <- ",
    "\u2194": " <-> ",
    # Superscripts / subscripts
    "\u00b2": "^2",
    "\u00b3": "^3",
    "\u2070": "^0",
    "\u00b9": "^1",
    "\u2074": "^4",
    "\u2075": "^5",
    "\u2076": "^6",
    "\u2077": "^7",
    "\u2078": "^8",
    "\u2079": "^9",
    "\u207f": "^n",
    "\u2080": "_0",
    "\u2081": "_1",
    "\u2082": "_2",
    "\u2083": "_3",
    "\u2084": "_4",
    # Misc
    "\u00d7": " * ",
    "\u00f7": " / ",
    "\u00b1": " +/- ",
    "\u2213": " -/+ ",
    "\u2026": "...",
}

# Characters considered "math" for density calculation
_MATH_CHARS = set(MATH_SYMBOL_MAP.keys()) | set("=+-*/^_{}[]()0123456789<>|\\")


@dataclass
class DocumentChunk:
    """A chunk of text with metadata for vector storage."""

    text: str
    embedding_text: str = ""
    metadata: dict = field(default_factory=dict)

    @property
    def chunk_id(self) -> str:
        """Unique ID based on source file + position."""
        source = self.metadata.get("source", "unknown")
        idx = self.metadata.get("chunk_index", 0)
        return f"{source}::chunk_{idx}"


class DocumentProcessor:
    """Extract text from documents and split into chunks."""

    def __init__(self):
        self.chunk_size = config.chunk_size
        self.chunk_overlap = config.chunk_overlap

    # ─── Main Entry Point ────────────────────────────────────────────────

    def process_file(
        self,
        file_path: Path,
        course_name: str = "",
        section_name: str = "",
        module_name: str = "",
    ) -> list[DocumentChunk]:
        """
        Extract text from a file, then chunk it.
        Returns list of DocumentChunks with rich metadata.
        """
        start = time.perf_counter()
        ext = file_path.suffix.lower()
        extractor = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".doc": self._extract_docx,  # best effort
            ".pptx": self._extract_pptx,
            ".txt": self._extract_text,
            ".md": self._extract_text,
            ".html": self._extract_html,
            ".htm": self._extract_html,
            ".rtf": self._extract_text,
        }.get(ext)

        if not extractor:
            logger.warning(f"Unsupported file format: {ext} ({file_path.name})")
            return []

        try:
            pages = extractor(file_path)
        except Exception as e:
            logger.error(f"Extraction failed [{file_path.name}]: {e}")
            return []

        if not pages:
            logger.warning(f"No text extracted from {file_path.name}")
            return []

        # Build metadata template
        base_meta = {
            "source": str(file_path),
            "filename": file_path.name,
            "course": course_name,
            "section": section_name,
            "module": module_name,
            "file_type": ext.lstrip("."),
        }

        # Chunk the pages
        chunks = self._chunk_pages(pages, base_meta)
        elapsed_ms = (time.perf_counter() - start) * 1000.0
        logger.info(f"Processed {file_path.name}: {len(pages)} pages → {len(chunks)} chunks in {elapsed_ms:.2f} ms")
        return chunks

    # ─── PDF Extraction ───────────────────────────────────────────────────

    def _extract_pdf(self, path: Path) -> list[str]:
        """Hybrid PDF extraction: pre-scans pages, routes text→pymupdf4llm, scanned→OCR.

        1. Quick pre-scan with fitz to classify each page (text vs scanned)
        2. Text-rich pages → pymupdf4llm batch processing (structured markdown)
        3. Scanned pages → direct OCR (skips pymupdf4llm entirely)
        No timeout needed. No quality loss. Fast on any PDF size.
        """
        try:
            import fitz
            import pymupdf4llm

            doc = fitz.open(str(path))
            total_pages = len(doc)

            # Phase 1: Pre-scan — classify pages as text or scanned
            text_page_indices = []
            scan_page_indices = []
            for i in range(total_pages):
                raw = doc[i].get_text("text").strip()
                if len(raw) > 100:
                    text_page_indices.append(i)
                else:
                    scan_page_indices.append(i)

            result_pages = {}  # page_index → extracted text
            ocr_count = 0

            # Phase 2: Scanned pages → OCR with early exit (majority vote)
            # Probe first 3 scanned pages; if 2+ fail quality check, skip the rest
            if scan_page_indices and OCR_AVAILABLE:
                PROBE_COUNT = min(3, len(scan_page_indices))

                # Phase 2a: Probe
                probe_fails = 0
                for idx in range(PROBE_COUNT):
                    pi = scan_page_indices[idx]
                    ocr_text = self._ocr_page(doc[pi], pi, path.name)
                    if ocr_text and self._ocr_quality_ok(ocr_text):
                        result_pages[pi] = ocr_text
                        ocr_count += 1
                    else:
                        probe_fails += 1

                # Phase 2b: Decide — majority failed → skip remaining
                remaining = len(scan_page_indices) - PROBE_COUNT
                if probe_fails >= 2 and remaining > 0:
                    logger.info(
                        f"OCR early exit: {probe_fails}/{PROBE_COUNT} probe pages unreadable, "
                        f"skipping {remaining} remaining in {path.name}"
                    )
                elif remaining > 0:
                    # Continue with remaining scanned pages
                    for idx in range(PROBE_COUNT, len(scan_page_indices)):
                        pi = scan_page_indices[idx]
                        ocr_text = self._ocr_page(doc[pi], pi, path.name)
                        if ocr_text and self._ocr_quality_ok(ocr_text):
                            result_pages[pi] = ocr_text
                            ocr_count += 1

                if ocr_count:
                    logger.info(f"OCR: {ocr_count}/{len(scan_page_indices)} scanned pages in {path.name}")

            doc.close()

            # Phase 3: Text pages → pymupdf4llm batch processing
            if text_page_indices:
                BATCH_SIZE = 50
                for batch_start in range(0, len(text_page_indices), BATCH_SIZE):
                    batch = text_page_indices[batch_start : batch_start + BATCH_SIZE]
                    try:
                        batch_data = pymupdf4llm.to_markdown(
                            str(path),
                            pages=batch,
                            page_chunks=True,
                            write_images=False,
                            show_progress=False,
                        )
                        for i, pd in enumerate(batch_data):
                            text = pd.get("text", "").strip()
                            if text:
                                result_pages[batch[i]] = text
                    except Exception as e:
                        logger.debug(f"pymupdf4llm batch failed for {path.name}: {e}")
                        doc = fitz.open(str(path))
                        for pi in batch:
                            text = doc[pi].get_text("text").strip()
                            if text:
                                result_pages[pi] = text
                        doc.close()

            # Assemble in page order
            pages = [result_pages[i] for i in sorted(result_pages.keys())]

            if pages:
                logger.info(
                    f"Extracted {len(pages)}p from {path.name} "
                    f"(text={len(text_page_indices)}, scanned={len(scan_page_indices)}, ocr={ocr_count})"
                )
                return pages

        except ImportError:
            logger.debug("pymupdf4llm not installed, using raw extraction")
        except Exception as e:
            logger.debug(f"pymupdf4llm failed for {path.name}: {e}")

        # Fallback: raw PyMuPDF extraction
        return self._extract_pdf_raw(path)

    def _extract_pdf_raw(self, path: Path) -> list[str]:
        """Raw PDF extraction using PyMuPDF with OCR fallback (legacy method)."""
        pages = []
        ocr_pages = 0

        try:
            import fitz

            doc = fitz.open(str(path))
            for page_num, page in enumerate(doc):
                text = page.get_text("text").strip()

                if len(text) < 50 and OCR_AVAILABLE:
                    ocr_text = self._ocr_page(page, page_num, path.name)
                    if ocr_text:
                        text = ocr_text
                        ocr_pages += 1

                if text:
                    pages.append(text)
            doc.close()

            if ocr_pages > 0:
                logger.info(f"OCR applied: {ocr_pages} page(s) in {path.name}")

            if pages:
                return pages
        except ImportError:
            pass
        except Exception as e:
            logger.debug(f"PyMuPDF failed, trying PyPDF2: {e}")

        # Fallback: PyPDF2 (no OCR support here)
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(str(path))
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(text.strip())
        except Exception as e:
            logger.error(f"PDF extraction failed [{path.name}]: {e}")

        return pages

    @staticmethod
    def _ocr_quality_ok(text: str) -> bool:
        """Check if OCR output is usable (not garbled noise).

        Uses ratio of alphabetic words (5+ chars) to total words.
        Real Turkish text: ~0.40-0.70 ratio, 100+ words.
        Garbled OCR (e.g. Ottoman manuscript): ~0.00-0.18 ratio.
        Threshold: ratio >= 0.25 AND count >= 15.
        """
        if len(text) < 50:
            return False
        words = text.split()
        if len(words) < 10:
            return False
        real_words = [w for w in words if w.isalpha() and len(w) >= 5]
        ratio = len(real_words) / len(words)
        return ratio >= 0.25 and len(real_words) >= 15

    def _ocr_page(self, page, page_num: int, filename: str) -> str:
        """OCR a single PDF page using Tesseract with two-pass strategy."""
        try:
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))

            # Pass 1: Turkish + English with PSM 6 (uniform block — better for academic pages)
            try:
                text = pytesseract.image_to_string(img, lang="tur+eng", config="--psm 6")
            except (OSError, RuntimeError, ValueError) as exc:
                logger.debug("OCR primary language pass failed on page %s: %s", page_num + 1, exc)
                text = pytesseract.image_to_string(img, lang="eng")

            # Pass 2: Extended languages + equ if pass 1 was insufficient
            if len(text.strip()) < 20:
                try:
                    text = pytesseract.image_to_string(img, lang="tur+eng+equ", config="--psm 6")
                except (OSError, RuntimeError, ValueError) as exc:
                    logger.debug("OCR secondary language pass failed on page %s: %s", page_num + 1, exc)
                    try:
                        text = pytesseract.image_to_string(img, lang="tur+eng+fra+deu+lat+ita+spa")
                        if text.strip():
                            logger.debug(f"OCR extended lang retry on page {page_num + 1} of {filename}")
                    except (OSError, RuntimeError, ValueError) as exc:
                        logger.debug("OCR extended language retry failed on page %s: %s", page_num + 1, exc)

            return text.strip()
        except Exception as e:
            logger.warning(f"OCR failed on page {page_num + 1} of {filename}: {e}")
            return ""

    # ─── Other Extractors ─────────────────────────────────────────────────

    def _extract_docx(self, path: Path) -> list[str]:
        """Extract text from DOCX."""
        try:
            from docx import Document

            doc = Document(str(path))
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [full_text] if full_text else []
        except Exception as e:
            logger.error(f"DOCX extraction failed [{path.name}]: {e}")
            return []

    def _extract_pptx(self, path: Path) -> list[str]:
        """Extract text from PPTX, one page per slide."""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
            return slides
        except Exception as e:
            logger.error(f"PPTX extraction failed [{path.name}]: {e}")
            return []

    def _extract_html(self, path: Path) -> list[str]:
        """Extract text from HTML."""
        try:
            from bs4 import BeautifulSoup

            html = path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(html, "html.parser")

            # Remove script/style tags
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)
            return [text] if text else []
        except Exception as e:
            logger.error(f"HTML extraction failed [{path.name}]: {e}")
            return []

    def _extract_text(self, path: Path) -> list[str]:
        """Extract plain text."""
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
            return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Text extraction failed [{path.name}]: {e}")
            return []

    # ─── Math Normalization ───────────────────────────────────────────────

    @staticmethod
    def _normalize_math_text(text: str) -> str:
        """Convert Unicode math symbols to searchable text for embedding."""
        for symbol, replacement in MATH_SYMBOL_MAP.items():
            if symbol in text:
                text = text.replace(symbol, replacement)
        # Clean up multiple spaces
        text = re.sub(r"  +", " ", text)
        return text

    @staticmethod
    def _math_density(line: str) -> float:
        """Calculate fraction of math-related characters in a line."""
        stripped = line.strip()
        if not stripped:
            return 0.0
        math_count = sum(1 for c in stripped if c in _MATH_CHARS)
        return math_count / len(stripped)

    @staticmethod
    def _has_math_content(text: str) -> bool:
        """Check if text contains significant math content."""
        math_count = sum(1 for c in text if c in set(MATH_SYMBOL_MAP.keys()))
        return math_count >= 3

    def _protect_equation_blocks(self, text: str) -> str:
        """Wrap consecutive math-heavy lines with sentinel markers."""
        lines = text.split("\n")
        result = []
        in_math = False

        for line in lines:
            density = self._math_density(line)
            is_math_line = density > 0.3 and len(line.strip()) > 5

            if is_math_line and not in_math:
                result.append("<<<MATH_BLOCK>>>")
                in_math = True
            elif not is_math_line and in_math and line.strip():
                result.append("<<<END_MATH>>>")
                in_math = False

            result.append(line)

        if in_math:
            result.append("<<<END_MATH>>>")

        return "\n".join(result)

    # ─── Chunking ─────────────────────────────────────────────────────────

    def _chunk_pages(self, pages: list[str], base_meta: dict) -> list[DocumentChunk]:
        """
        Split extracted pages into overlapping chunks.
        Produces dual-text chunks: original for LLM, normalized for embedding.
        """
        # Combine all pages with page markers
        full_text = ""
        for i, page in enumerate(pages):
            full_text += f"\n[Page {i+1}]\n{page}\n"

        # Protect equation blocks from splitting
        protected_text = self._protect_equation_blocks(full_text)

        # Split with equation-aware separators
        chunks_text = self._recursive_split(protected_text)

        # Build dual-text chunks
        has_math = self._has_math_content(full_text)
        result = []
        for i, chunk_text in enumerate(chunks_text):
            # Remove sentinel markers from final text
            clean_text = chunk_text.replace("<<<MATH_BLOCK>>>", "").replace("<<<END_MATH>>>", "")
            clean_text = clean_text.strip()
            if not clean_text:
                continue

            # Normalized text for embedding
            embedding_text = self._normalize_math_text(clean_text)

            meta = {
                **base_meta,
                "chunk_index": i,
                "total_chunks": len(chunks_text),
                "has_math": has_math,
            }
            result.append(
                DocumentChunk(
                    text=clean_text,
                    embedding_text=embedding_text,
                    metadata=meta,
                )
            )

        return result

    def _recursive_split(self, text: str) -> list[str]:
        """Split text into chunks with overlap. Uses equation-aware separators."""
        try:
            from langchain_text_splitters import RecursiveCharacterTextSplitter

            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=[
                    "<<<END_MATH>>>\n<<<MATH_BLOCK>>>",  # between equation blocks
                    "<<<END_MATH>>>",  # after equation block
                    "<<<MATH_BLOCK>>>",  # before equation block
                    "\n\n",
                    "\n",
                    ". ",
                    " ",
                    "",
                ],
            )
            return splitter.split_text(text)
        except ImportError:
            # Manual fallback
            return self._manual_split(text)

    def _manual_split(self, text: str) -> list[str]:
        """Simple chunking fallback without langchain."""
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size

            # Try to break at sentence boundary
            if end < len(text):
                for sep in ["\n\n", "\n", ". ", " "]:
                    idx = text.rfind(sep, start, end)
                    if idx > start:
                        end = idx + len(sep)
                        break

            chunks.append(text[start:end])
            start = end - self.chunk_overlap

        return chunks
